from pathlib import Path

import fitz
from docx import Document
from pptx import Presentation


class MaterialParser:

    @staticmethod
    def parse(path: str) -> str:

        extension = Path(path).suffix.lower()

        if extension == ".pdf":
            return MaterialParser.parse_pdf(path)

        if extension == ".docx":
            return MaterialParser.parse_docx(path)

        if extension == ".pptx":
            return MaterialParser.parse_pptx(path)

        if extension == ".txt":
            with open(
                path,
                encoding="utf-8",
            ) as f:
                return f.read()

        return ""

    @staticmethod
    def parse_pdf(path: str) -> str:

        text = ""

        pdf = fitz.open(path)

        for page in pdf:
            text += page.get_text()

        pdf.close()

        return text

    @staticmethod
    def parse_docx(path: str) -> str:

        document = Document(path)

        return "\n".join(
            paragraph.text
            for paragraph in document.paragraphs
        )

    @staticmethod
    def parse_pptx(path: str) -> str:

        presentation = Presentation(path)

        text = ""

        for slide in presentation.slides:

            for shape in slide.shapes:

                if hasattr(shape, "text"):
                    text += shape.text + "\n"

        return text